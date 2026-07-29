"""Aplica o design system VAP (Vendas de Alta Performance) no deck da Man Motors.

Troca a paleta azul/vermelha pela paleta ouro/preto/creme, muda a tipografia para
Poppins, normaliza os paddings internos e substitui os logos pelo da Jocelaine.

Uso:  python aplicar-ds-vap.py entrada.pptx saida.pptx
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image

# ---------------------------------------------------------------- design system
INK       = "151210"   # preto
INK_SOFT  = "5F574E"
PAPER     = "FDFAF4"   # creme claro (fundo)
PAPER_2   = "F4ECDC"   # creme painel
GOLD      = "C08A2A"
GOLD_DEEP = "96651A"
GOLD_LT   = "E6BD5E"
LINE      = "E3D3B0"
WHITE     = "FFFFFF"

# tipografia herdada do sistema de design do Metodo LAP
FONT_TITLE = "Playfair Display"     # titulos e numerais
FONT_BODY  = "DM Sans 14pt"         # corpo, rotulos e rodape
FONT = FONT_BODY

# rodape: sai a marca do cliente, entra o nome do treinamento
FOOTER_OLD = "MAN MOTORS S.A."
FOOTER_NEW = "VENDAS DE ALTA PERFORMANCE  •  O CÓDIGO DAS VENDAS INFALÍVEIS"


def is_title(shape_name, run):
    """Titulos e numerais do LAP: negrito a partir de 16pt."""
    if not run.font.bold:
        return False
    sz = run.font.size.pt if run.font.size else 0
    if sz < 16:
        return False
    return not (shape_name == "Text 9" and sz == 17.0)   # apoio da capa fica no corpo


def set_spacing(run, hundredths):
    """letter-spacing do rodape (o python-pptx nao expoe; vai no XML)."""
    run.font._rPr.set("spc", str(hundredths))
# emoji colorido foge da paleta e nem toda maquina tem o glifo:
# vira um losango dourado, que existe em qualquer fonte
ICON = "◆"


def is_emoji(txt):
    t = txt.strip()
    return bool(t) and all(ord(ch) > 0x2100 for ch in t)

# cor antiga -> cor nova
FILL_MAP = {
    "0B2D6B": INK,        # azul estrutural (barras, faixas)
    "E30613": GOLD,       # vermelho de destaque
    "F4F7FB": PAPER_2,    # painel claro
    "CCD6E3": LINE,       # bordas
    "D4DEE9": LINE,
    "CBD5E1": LINE,
}
TEXT_MAP = {
    "0B2D6B": INK,
    "1F2937": INK,
    "4B5563": INK_SOFT,
    "E30613": GOLD_DEEP,  # destaque sobre fundo claro
    "BFE3FF": GOLD_LT,    # destaque sobre fundo escuro
}
RED_OLD = "E30613"

PAD = Inches(0.09)


def hexof(color):
    try:
        return str(color.rgb).upper()
    except Exception:
        return None


def shape_fill_hex(sh):
    try:
        if sh.fill.type == 1:
            return hexof(sh.fill.fore_color)
    except Exception:
        pass
    return None


def bbox(sh):
    if None in (sh.left, sh.top, sh.width, sh.height):
        return None
    return (sh.left, sh.top, sh.left + sh.width, sh.top + sh.height)


def center_inside(inner, outer):
    if not inner or not outer:
        return False
    cx = (inner[0] + inner[2]) / 2
    cy = (inner[1] + inner[3]) / 2
    return outer[0] <= cx <= outer[2] and outer[1] <= cy <= outer[3]


def make_logo(src, out, aspect, invert_to_white=False):
    """Encaixa o logo da Jocelaine numa tela com o aspecto do quadro original,
    para o PowerPoint nao esticar a marca."""
    im = Image.open(src).convert("RGBA")
    w = im.width
    h = int(round(w / aspect))
    if h < im.height:                      # tela mais baixa que o logo: alarga
        h = im.height
        w = int(round(h * aspect))
    canvas = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    canvas.alpha_composite(im, ((w - im.width) // 2, (h - im.height) // 2))
    canvas.save(out)
    return out


def restyle(src_path, out_path):
    prs = Presentation(src_path)

    # ---------------------------------------------------------- logos
    # todos os quadros de logo ficam sobre area escura -> versao branca
    logo = make_logo("logo-jr.png", "_logo_vap.png", 334 / 150)
    blob = open(logo, "rb").read()
    seen = set()
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.shape_type == 13:
                part = sh.image  # ImagePart wrapper
                ip = sh._element.blip_rId
                image_part = slide.part.related_part(ip)
                if id(image_part) not in seen:
                    image_part._blob = blob
                    image_part.blob  # noqa
                    seen.add(id(image_part))

    # python-pptx guarda o blob no _blob da part; forca a regravacao
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.shape_type == 13:
                image_part = slide.part.related_part(sh._element.blip_rId)
                image_part._blob = blob

    for idx, slide in enumerate(prs.slides, 1):
        # ------------------------------------------------------ fundo creme
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor.from_string(INK if idx == 13 else PAPER)

        # guarda quem era vermelho, para corrigir o texto branco em cima
        red_boxes = [bbox(sh) for sh in slide.shapes if shape_fill_hex(sh) == RED_OLD]

        for sh in slide.shapes:
            # -------------------------------------------------- rodape
            if sh.has_text_frame and FOOTER_OLD in sh.text_frame.text.upper():
                for para in sh.text_frame.paragraphs:
                    if FOOTER_OLD not in "".join(r.text for r in para.runs).upper():
                        continue
                    for n, r in enumerate(para.runs):
                        r.text = FOOTER_NEW if n == 0 else ""

            # -------------------------------------------------- preenchimentos
            cur = shape_fill_hex(sh)
            if cur in FILL_MAP:
                sh.fill.fore_color.rgb = RGBColor.from_string(FILL_MAP[cur])
            try:
                lc = hexof(sh.line.color) if sh.line.fill.type == 1 else None
                if lc in FILL_MAP:
                    sh.line.color.rgb = RGBColor.from_string(FILL_MAP[lc])
            except Exception:
                pass

            # -------------------------------------------------- tipografia
            if sh.has_text_frame:
                tf = sh.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_right = PAD
                tf.margin_top = tf.margin_bottom = Emu(int(PAD * 0.5))
                on_gold = any(center_inside(bbox(sh), rb) for rb in red_boxes)
                for para in tf.paragraphs:
                    for r in para.runs:
                        if is_emoji(r.text):
                            r.text = ICON
                            r.font.name = FONT
                            r.font.color.rgb = RGBColor.from_string(GOLD)
                            continue
                        r.font.name = FONT_TITLE if is_title(sh.name, r) else FONT_BODY
                        if r.font.size and r.font.size.pt <= 11:
                            set_spacing(r, 60)      # rotulos e rodape respiram
                        c = hexof(r.font.color)
                        if c in TEXT_MAP:
                            r.font.color.rgb = RGBColor.from_string(TEXT_MAP[c])
                        elif c == WHITE and on_gold:
                            # branco sobre ouro nao tem contraste -> preto
                            r.font.color.rgb = RGBColor.from_string(INK)

        # ------------------------------------------------------ paddings dos logos
        for sh in slide.shapes:
            if sh.shape_type != 13:
                continue
            if idx == 1:                      # capa: centraliza no painel escuro
                sh.width = Inches(3.30)
                sh.height = Inches(3.30 * 150 / 334)
                sh.left = Inches((4.70 - 3.30) / 2)
                sh.top = int((prs.slide_height - sh.height) / 2)
            elif idx == 13:                   # encerramento
                sh.width = Inches(5.60)
                sh.height = Inches(5.60 * 150 / 334)
                sh.left = int((prs.slide_width - sh.width) / 2)
                sh.top = int((prs.slide_height - sh.height) / 2)
            else:                             # cabecalho: dentro da barra preta
                sh.height = Inches(0.52)
                sh.width = Inches(0.52 * 334 / 150)
                sh.left = Inches(0.38)
                sh.top = Inches(0.115)

    # ---------------------------------------------------------- capa
    s1 = prs.slides[0]
    for sh in s1.shapes:
        if sh.name == "Shape 1":              # painel lateral -> preto
            sh.fill.fore_color.rgb = RGBColor.from_string(INK)
            sh.line.color.rgb = RGBColor.from_string(INK)

    # ---------------------------------------------------------- encerramento
    s13 = prs.slides[12]
    for sh in s13.shapes:
        if sh.name == "Shape 0":              # cobre a pagina inteira
            sh.left, sh.top = 0, 0
            sh.width, sh.height = prs.slide_width, prs.slide_height
            sh.fill.fore_color.rgb = RGBColor.from_string(INK)
            sh.line.color.rgb = RGBColor.from_string(INK)

    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    src, out = sys.argv[1], sys.argv[2]
    print("gerado:", restyle(src, out))
